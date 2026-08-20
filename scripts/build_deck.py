"""Generate SKILLBRIDGE hackathon pitch deck as a .pptx file."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree
import qrcode
from PIL import Image
import io, os

DEMO_URL = "https://ai-match-hub-10.preview.emergentagent.com/"
QR_PATH = "/app/scripts/_qr_demo.png"

def build_qr(url, path):
    qr = qrcode.QRCode(version=None, error_correction=qrcode.constants.ERROR_CORRECT_H, box_size=20, border=2)
    qr.add_data(url); qr.make(fit=True)
    img = qr.make_image(fill_color=(8, 10, 18), back_color=(246, 247, 251)).convert("RGBA")
    img.save(path)
    return path

build_qr(DEMO_URL, QR_PATH)

# ---------- Brand palette ----------
BG = RGBColor(0x08, 0x0A, 0x12)           # near-black
SURFACE = RGBColor(0x14, 0x1A, 0x2B)       # card
SURFACE2 = RGBColor(0x0E, 0x12, 0x20)      # deeper card
LINE = RGBColor(0x2A, 0x33, 0x50)
MUTED = RGBColor(0x92, 0x99, 0xAD)
WHITE = RGBColor(0xF6, 0xF7, 0xFB)
CYAN = RGBColor(0x66, 0xE4, 0xEE)
VIOLET = RGBColor(0x9A, 0x7A, 0xFF)
PINK = RGBColor(0xF7, 0x8A, 0xD5)
AMBER = RGBColor(0xF4, 0xBC, 0x6B)
GREEN = RGBColor(0x4F, 0xD4, 0x8A)

# 16:9 widescreen
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

BLANK = prs.slide_layouts[6]


def add_bg(slide, color=BG):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.shadow.inherit = False
    return bg


def add_text(slide, left, top, width, height, text, *, size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT, font="Calibri", anchor=MSO_ANCHOR.TOP, line_spacing=1.15, letter_spacing=None):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        if letter_spacing is not None:
            rPr = r._r.get_or_add_rPr()
            rPr.set("spc", str(letter_spacing))
    return tb


def add_eyebrow(slide, left, top, text, color=CYAN, width=Inches(6)):
    return add_text(slide, left, top, width, Inches(0.3), text.upper(), size=10, bold=True, color=color, font="Consolas", letter_spacing=180)


def add_rect(slide, left, top, width, height, fill=SURFACE, line=LINE, corner=None):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill
    rect.line.color.rgb = line
    rect.line.width = Pt(0.75)
    rect.shadow.inherit = False
    return rect


def add_pill(slide, left, top, text, color=CYAN, bg=None, width=None):
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width or Inches(1.8), Inches(0.32))
    pill.fill.solid()
    pill.fill.fore_color.rgb = bg or RGBColor(0x10, 0x1F, 0x27)
    pill.line.color.rgb = color
    pill.line.width = Pt(0.75)
    pill.shadow.inherit = False
    tf = pill.text_frame
    tf.margin_left = tf.margin_right = Inches(0.1)
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = "Consolas"
    r.font.size = Pt(9)
    r.font.bold = True
    r.font.color.rgb = color
    return pill


def add_grid_bg(slide, size=Inches(0.7)):
    # subtle grid lines
    n_x = int(SW / size) + 1
    n_y = int(SH / size) + 1
    for i in range(1, n_x):
        ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, i * size, 0, Emu(4000), SH)
        ln.fill.solid()
        ln.fill.fore_color.rgb = RGBColor(0x14, 0x18, 0x28)
        ln.line.fill.background()
        ln.shadow.inherit = False
    for i in range(1, n_y):
        ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, i * size, SW, Emu(4000))
        ln.fill.solid()
        ln.fill.fore_color.rgb = RGBColor(0x14, 0x18, 0x28)
        ln.line.fill.background()
        ln.shadow.inherit = False


def add_brand_mark(slide, left, top, size=Inches(0.45)):
    diamond = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, left, top, size, size)
    diamond.fill.background()
    diamond.line.color.rgb = CYAN
    diamond.line.width = Pt(1.4)
    diamond.shadow.inherit = False
    inner_size = size / 2.3
    inner = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + size / 2 - inner_size / 2, top + size / 2 - inner_size / 2, inner_size, inner_size)
    inner.fill.solid()
    inner.fill.fore_color.rgb = VIOLET
    inner.line.fill.background()
    inner.shadow.inherit = False


def add_brand_lockup(slide, left, top):
    add_brand_mark(slide, left, top)
    tb = add_text(slide, left + Inches(0.6), top - Inches(0.02), Inches(3), Inches(0.5), "", size=13, bold=True, color=WHITE, letter_spacing=180)
    tf = tb.text_frame
    p = tf.paragraphs[0]
    r1 = p.add_run(); r1.text = "SKILL"; r1.font.name = "Calibri"; r1.font.size = Pt(13); r1.font.bold = True; r1.font.color.rgb = WHITE
    rPr = r1._r.get_or_add_rPr(); rPr.set("spc", "220")
    r2 = p.add_run(); r2.text = "BRIDGE"; r2.font.name = "Calibri"; r2.font.size = Pt(13); r2.font.bold = True; r2.font.color.rgb = CYAN
    rPr = r2._r.get_or_add_rPr(); rPr.set("spc", "220")


def slide_number(slide, n, total):
    add_text(slide, SW - Inches(1.2), SH - Inches(0.5), Inches(1), Inches(0.3), f"{n:02d} / {total:02d}", size=10, color=MUTED, font="Consolas", align=PP_ALIGN.RIGHT, letter_spacing=120)


def footer_line(slide):
    add_text(slide, Inches(0.6), SH - Inches(0.5), Inches(8), Inches(0.3), "SKILLBRIDGE  ·  BRIDGE YOUR SKILLS. DISCOVER YOUR OPPORTUNITIES.", size=9, color=MUTED, font="Consolas", letter_spacing=180)


TOTAL = 12
n = 1


def new_slide(title_eyebrow=None):
    global n
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_grid_bg(s)
    add_brand_lockup(s, Inches(0.6), Inches(0.45))
    if title_eyebrow:
        add_eyebrow(s, Inches(0.6), Inches(1.2), title_eyebrow)
    slide_number(s, n, TOTAL)
    footer_line(s)
    n += 1
    return s


# =====================================================================
# 1. TITLE SLIDE
# =====================================================================
s1 = prs.slides.add_slide(BLANK)
add_bg(s1)
add_grid_bg(s1, size=Inches(0.85))

# soft glow blobs
for cx, cy, rgb in [(Inches(9.5), Inches(2.2), VIOLET), (Inches(10.5), Inches(5), CYAN)]:
    blob = s1.shapes.add_shape(MSO_SHAPE.OVAL, cx, cy, Inches(3.6), Inches(3.6))
    blob.fill.solid()
    blob.fill.fore_color.rgb = rgb
    blob.line.fill.background()
    blob.fill.transparency = 0.85
    blob.shadow.inherit = False

add_brand_lockup(s1, Inches(0.6), Inches(0.45))
add_eyebrow(s1, Inches(11), Inches(0.55), "· AI HACKATHON", color=CYAN, width=Inches(2))

add_eyebrow(s1, Inches(0.6), Inches(2.4), "BRIDGE YOUR NEXT CHAPTER", color=CYAN)

# Main headline
tb = s1.shapes.add_textbox(Inches(0.6), Inches(2.75), Inches(12), Inches(3.3))
tf = tb.text_frame
tf.word_wrap = True
tf.margin_left = tf.margin_right = 0
p1 = tf.paragraphs[0]
p1.line_spacing = 1.02
r = p1.add_run(); r.text = "Bridge your skills."; r.font.name = "Calibri"; r.font.size = Pt(72); r.font.bold = True; r.font.color.rgb = WHITE
p2 = tf.add_paragraph()
p2.line_spacing = 1.02
r = p2.add_run(); r.text = "Discover "; r.font.name = "Calibri"; r.font.size = Pt(72); r.font.bold = True; r.font.color.rgb = VIOLET
r = p2.add_run(); r.text = "your opportunities."; r.font.name = "Calibri"; r.font.size = Pt(72); r.font.bold = True; r.font.color.rgb = WHITE

add_text(s1, Inches(0.6), Inches(6.0), Inches(9), Inches(0.5), "AI-powered opportunities built around what you can do.", size=18, color=MUTED)

# hero pill row
add_pill(s1, Inches(0.6), Inches(6.7), "AI SKILL MATCH", CYAN, width=Inches(1.9))
add_pill(s1, Inches(2.7), Inches(6.7), "STUDENTS + PROS", VIOLET, width=Inches(2.1))
add_pill(s1, Inches(5.0), Inches(6.7), "HIRING PIPELINE", PINK, width=Inches(2.1))

add_text(s1, Inches(9.5), Inches(6.7), Inches(3.3), Inches(0.4), "Team · Hackathon 2026", size=11, color=MUTED, font="Consolas", align=PP_ALIGN.RIGHT, letter_spacing=180)


# =====================================================================
# 2. THE PROBLEM
# =====================================================================
s2 = new_slide("01 · THE PROBLEM")
add_text(s2, Inches(0.6), Inches(1.55), Inches(12), Inches(1.8), "Talent and opportunity\nkeep missing each other.", size=54, bold=True, color=WHITE, line_spacing=1.02)

problems = [
    ("Generic listings", "Students scroll past hundreds of roles they don’t fit.\nRelevance is buried inside filters that need a resume they don’t have yet.", VIOLET),
    ("Skills go invisible", "Freelancers and skilled professionals waste hours proving what a\n2-line search bar cannot see — projects, side wins, portfolios.", CYAN),
    ("Hiring is fragmented", "Owners of internships and freelance briefs juggle inboxes, forms and\nspreadsheets — no clean pipeline from applicant to shortlist.", PINK),
]
top = Inches(4.2)
card_w = Inches(4.05)
gap = Inches(0.15)
for i, (title, body, color) in enumerate(problems):
    left = Inches(0.6) + i * (card_w + gap)
    add_rect(s2, left, top, card_w, Inches(2.5))
    # top accent
    acc = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, card_w, Inches(0.05))
    acc.fill.solid(); acc.fill.fore_color.rgb = color; acc.line.fill.background(); acc.shadow.inherit = False
    add_text(s2, left + Inches(0.3), top + Inches(0.35), card_w - Inches(0.6), Inches(0.5), title, size=18, bold=True, color=WHITE)
    add_text(s2, left + Inches(0.3), top + Inches(0.9), card_w - Inches(0.6), Inches(1.6), body, size=12, color=MUTED, line_spacing=1.4)


# =====================================================================
# 3. THE SOLUTION
# =====================================================================
s3 = new_slide("02 · THE SOLUTION")
add_text(s3, Inches(0.6), Inches(1.55), Inches(9), Inches(1.6), "One AI-powered bridge.\nBoth sides of the market.", size=48, bold=True, color=WHITE, line_spacing=1.05)

add_text(s3, Inches(0.6), Inches(3.7), Inches(7.2), Inches(2.5),
         "SKILLBRIDGE learns what a person can actually do — skills, interests,\n"
         "portfolio, past applications — and matches them to the opportunities\n"
         "where they will meaningfully show up.\n\n"
         "For students it removes the noise. For professionals it surfaces work\n"
         "worth their craft. For owners it turns an inbox into a pipeline.",
         size=15, color=MUTED, line_spacing=1.55)

# Diagram: 3 nodes with arrows to central AI core
core_left, core_top = Inches(9.6), Inches(4.05)
core = s3.shapes.add_shape(MSO_SHAPE.OVAL, core_left, core_top, Inches(2.4), Inches(2.4))
core.fill.solid(); core.fill.fore_color.rgb = SURFACE
core.line.color.rgb = VIOLET; core.line.width = Pt(1.5)
core.shadow.inherit = False
# glow rings
for r_off, alpha in [(Inches(0.35), 0.82), (Inches(0.7), 0.9)]:
    ring = s3.shapes.add_shape(MSO_SHAPE.OVAL, core_left - r_off, core_top - r_off, Inches(2.4) + r_off * 2, Inches(2.4) + r_off * 2)
    ring.fill.background()
    ring.line.color.rgb = CYAN
    ring.line.width = Pt(0.5)
    ring.line.dash_style = 7  # dash
    ring.shadow.inherit = False

add_text(s3, core_left, core_top + Inches(0.75), Inches(2.4), Inches(0.5), "AI SKILL MATCH", size=10, bold=True, color=CYAN, font="Consolas", align=PP_ALIGN.CENTER, letter_spacing=180)
add_text(s3, core_left, core_top + Inches(1.05), Inches(2.4), Inches(0.7), "89%", size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# =====================================================================
# 4. KEY FEATURES
# =====================================================================
s4 = new_slide("03 · WHAT WE BUILT")
add_text(s4, Inches(0.6), Inches(1.55), Inches(12), Inches(1.4), "Everything you need for a\nfull two-sided demo.", size=44, bold=True, color=WHITE, line_spacing=1.05)

features = [
    ("Role Selection", "Student / Skilled Professional flow with tailored onboarding.", VIOLET),
    ("AI Skill Match", "Deterministic overlap scoring · Home + Dashboard + Skill Match page.", CYAN),
    ("Personalised Feed", "Opportunities ranked by user skills, interests and goals.", PINK),
    ("Application Flow", "Multi-step Apply · Resume · Cover letter · Instant success.", AMBER),
    ("Applicant Pipeline", "Owners review applicants, move Applied → Shortlisted → Selected.", GREEN),
    ("Portfolio Uploads", "PDF / images via Emergent Object Storage — served through API.", VIOLET),
    ("Notification Bell", "Live popover · Unread dot · Mark-all-read · Deep link inbox.", CYAN),
    ("Real Reset Email", "Resend integration · Branded HTML · Sandbox-aware fallback.", PINK),
]
top = Inches(3.55)
cols = 4
card_w = (SW - Inches(0.6) * 2 - Inches(0.2) * (cols - 1)) / cols
card_h = Inches(1.5)
for i, (title, body, color) in enumerate(features):
    row = i // cols
    col = i % cols
    left = Inches(0.6) + col * (card_w + Inches(0.2))
    t = top + row * (card_h + Inches(0.2))
    add_rect(s4, left, t, card_w, card_h)
    dot = s4.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.25), t + Inches(0.28), Inches(0.16), Inches(0.16))
    dot.fill.solid(); dot.fill.fore_color.rgb = color; dot.line.fill.background(); dot.shadow.inherit = False
    add_text(s4, left + Inches(0.55), t + Inches(0.22), card_w - Inches(0.7), Inches(0.35), title, size=14, bold=True, color=WHITE)
    add_text(s4, left + Inches(0.25), t + Inches(0.7), card_w - Inches(0.5), Inches(0.9), body, size=10.5, color=MUTED, line_spacing=1.4)


# =====================================================================
# 5. AI SKILL MATCH DEEP DIVE
# =====================================================================
s5 = new_slide("04 · THE AI MATCH")
add_text(s5, Inches(0.6), Inches(1.55), Inches(9), Inches(1.6), "Deterministic, explainable,\nfast enough to feel magical.", size=42, bold=True, color=WHITE, line_spacing=1.05)

add_text(s5, Inches(0.6), Inches(3.6), Inches(7), Inches(2.5),
         "Same input → same score. No black-box tricks.\n\n"
         "•  Overlap ratio between user skills and role skills\n"
         "•  +10 pt bonus when 2+ skills align (signal, not luck)\n"
         "•  Boosts for interest + goal alignment\n"
         "•  Clamped 38 – 99 so no score feels fake\n\n"
         "Every match on Home, Dashboard and Detail is computed live from\n"
         "the user’s profile. Edit a skill → the whole app re-personalises.",
         size=14, color=MUTED, line_spacing=1.5)

# Right side: mock skill match card
card_left = Inches(8.6)
card_top = Inches(3.3)
add_rect(s5, card_left, card_top, Inches(4.2), Inches(3.3), fill=SURFACE, line=CYAN)
add_eyebrow(s5, card_left + Inches(0.3), card_top + Inches(0.3), "AI SKILL MATCH", color=CYAN)
add_text(s5, card_left + Inches(0.3), card_top + Inches(0.55), Inches(3.6), Inches(1.1), "89%", size=64, bold=True, color=WHITE)

# score bars
skills = [("React", 0.94, GREEN), ("JavaScript", 0.88, GREEN), ("UI/UX", 0.76, AMBER), ("TypeScript", 0.42, VIOLET)]
bar_top = card_top + Inches(1.85)
for i, (name, val, color) in enumerate(skills):
    y = bar_top + i * Inches(0.32)
    add_text(s5, card_left + Inches(0.3), y, Inches(1.6), Inches(0.25), name, size=11, color=WHITE)
    add_text(s5, card_left + Inches(3.4), y, Inches(0.7), Inches(0.25), f"{int(val*100)}%", size=10, color=CYAN, font="Consolas", align=PP_ALIGN.RIGHT)
    # bar
    track = s5.shapes.add_shape(MSO_SHAPE.RECTANGLE, card_left + Inches(1.5), y + Inches(0.05), Inches(1.8), Inches(0.13))
    track.fill.solid(); track.fill.fore_color.rgb = RGBColor(0x24, 0x29, 0x3A); track.line.fill.background(); track.shadow.inherit = False
    fill = s5.shapes.add_shape(MSO_SHAPE.RECTANGLE, card_left + Inches(1.5), y + Inches(0.05), Inches(1.8) * val, Inches(0.13))
    fill.fill.solid(); fill.fill.fore_color.rgb = color; fill.line.fill.background(); fill.shadow.inherit = False


# =====================================================================
# 6. TWO-SIDED DEMO FLOW
# =====================================================================
s6 = new_slide("05 · THE DEMO JOURNEY")
add_text(s6, Inches(0.6), Inches(1.55), Inches(12), Inches(1.4), "A full loop in under two minutes.", size=40, bold=True, color=WHITE)

# Two lanes
lanes = [
    ("STUDENT · AARAV", VIOLET, [
        "Role selection · Continue as Student",
        "Demo access · Onboarding (skills, interests, goal)",
        "Home shows personalised match + top 3 opportunities",
        "Open Frontend Dev Intern · Apply flow · Submit",
        "Application appears in My Applications instantly",
    ]),
    ("PROFESSIONAL · MAYA (OWNER)", CYAN, [
        "Enter demo as Skilled Professional",
        "Tune profile · Upload portfolio PDF to storage",
        "Workspace shows 7 owned opportunities with counts",
        "Review Applicants · See Aarav ranked by match score",
        "Move to Shortlisted → Aarav gets a live notification",
    ]),
]
top = Inches(3.15)
lane_w = Inches(6.05)
for i, (title, color, steps) in enumerate(lanes):
    left = Inches(0.6) + i * (lane_w + Inches(0.2))
    add_rect(s6, left, top, lane_w, Inches(3.6))
    acc = s6.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.06), Inches(3.6))
    acc.fill.solid(); acc.fill.fore_color.rgb = color; acc.line.fill.background(); acc.shadow.inherit = False
    add_eyebrow(s6, left + Inches(0.35), top + Inches(0.25), title, color=color, width=lane_w)
    for j, step in enumerate(steps):
        y = top + Inches(0.75) + j * Inches(0.55)
        num = s6.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.35), y, Inches(0.32), Inches(0.32))
        num.fill.background(); num.line.color.rgb = color; num.line.width = Pt(1); num.shadow.inherit = False
        add_text(s6, left + Inches(0.35), y - Inches(0.02), Inches(0.32), Inches(0.32), str(j + 1), size=11, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text(s6, left + Inches(0.85), y + Inches(0.03), lane_w - Inches(1.05), Inches(0.5), step, size=12, color=WHITE, line_spacing=1.4)


# =====================================================================
# 7. TECH ARCHITECTURE
# =====================================================================
s7 = new_slide("06 · ARCHITECTURE")
add_text(s7, Inches(0.6), Inches(1.55), Inches(12), Inches(1.4), "Small stack. Big product.", size=40, bold=True, color=WHITE)

# Three columns: Frontend / Backend / Integrations
tiers = [
    ("FRONTEND", VIOLET, "React 19 + React Router\nContext-driven dynamic data\nDeterministic match (mirror)\nSkeletons + empty states\nTailwind-free custom CSS"),
    ("BACKEND", CYAN, "FastAPI + Motor (async)\nMongoDB collections:\n · users · opportunities\n · applications · saved\n · notifications · portfolio\nSeeded opportunities on startup"),
    ("INTEGRATIONS", PINK, "Resend — real reset email\nEmergent Object Storage —\n   portfolio uploads (PDF / img)\nPydantic v2 payload models\nCORS + supervisor-managed"),
]
# three tier cards
top = Inches(3.15)
col_w = Inches(4.05)
gap = Inches(0.15)
for i, (title, color, body) in enumerate(tiers):
    left = Inches(0.6) + i * (col_w + gap)
    add_rect(s7, left, top, col_w, Inches(3.1))
    add_eyebrow(s7, left + Inches(0.3), top + Inches(0.3), title, color=color, width=col_w)
    add_text(s7, left + Inches(0.3), top + Inches(0.7), col_w - Inches(0.6), Inches(2.3), body, size=12.5, color=MUTED, line_spacing=1.5)

# small stats row
stats = [("18", "Pytest cases pass"), ("10", "Seeded opportunities"), ("6", "Application statuses"), ("0", "Console errors")]
top = Inches(6.4)
w = Inches(3.02)
for i, (v, l) in enumerate(stats):
    left = Inches(0.6) + i * (w + Inches(0.12))
    add_rect(s7, left, top, w, Inches(0.5), fill=SURFACE2)
    add_text(s7, left + Inches(0.2), top + Inches(0.09), Inches(0.7), Inches(0.4), v, size=16, bold=True, color=CYAN)
    add_text(s7, left + Inches(0.95), top + Inches(0.15), w - Inches(1.05), Inches(0.4), l, size=9.5, color=MUTED, font="Consolas", letter_spacing=140)


# =====================================================================
# 8. KEY SCREENS (visual placeholders)
# =====================================================================
s8 = new_slide("07 · THE PRODUCT · SCREENS")
add_text(s8, Inches(0.6), Inches(1.55), Inches(12), Inches(1.4), "Real, working, live.", size=40, bold=True, color=WHITE)

screens = [
    ("HOME · MATCH", "Dynamic score · strongest skill · next skill\nPersonalised top 3 opportunities", VIOLET),
    ("APPLY FLOW", "4 steps · resume upload · cover letter\nSuccess screen · applications updates live", CYAN),
    ("APPLICANT PIPELINE", "Ranked by match · portfolio link\nStatus buttons write to DB + notify user", PINK),
    ("NOTIFICATION BELL", "Popover · unread dot · mark-all-read\nStatus changes appear in real time", AMBER),
]
top = Inches(3.15)
w = (SW - Inches(0.6) * 2 - Inches(0.15) * 3) / 4
for i, (t, body, color) in enumerate(screens):
    left = Inches(0.6) + i * (w + Inches(0.15))
    add_rect(s8, left, top, w, Inches(3.5))
    # mock screen top bar
    bar = s8.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + Inches(0.2), top + Inches(0.25), w - Inches(0.4), Inches(0.4))
    bar.fill.solid(); bar.fill.fore_color.rgb = SURFACE2; bar.line.fill.background(); bar.shadow.inherit = False
    # mock diagonal graphic
    stripe = s8.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + Inches(0.2), top + Inches(0.8), w - Inches(0.4), Inches(1.4))
    stripe.fill.solid(); stripe.fill.fore_color.rgb = SURFACE2; stripe.line.fill.background(); stripe.shadow.inherit = False
    add_text(s8, left + Inches(0.25), top + Inches(0.95), w - Inches(0.5), Inches(0.4), t, size=13, bold=True, color=color)
    # progress mock
    for j in range(3):
        y = top + Inches(1.5) + j * Inches(0.22)
        track = s8.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + Inches(0.35), y, w - Inches(0.7), Inches(0.09))
        track.fill.solid(); track.fill.fore_color.rgb = RGBColor(0x24, 0x29, 0x3A); track.line.fill.background(); track.shadow.inherit = False
        fill = s8.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + Inches(0.35), y, (w - Inches(0.7)) * (0.9 - j * 0.15), Inches(0.09))
        fill.fill.solid(); fill.fill.fore_color.rgb = color; fill.line.fill.background(); fill.shadow.inherit = False
    add_text(s8, left + Inches(0.25), top + Inches(2.4), w - Inches(0.5), Inches(1), body, size=11, color=MUTED, line_spacing=1.45)


# =====================================================================
# 9. WHAT MAKES US DIFFERENT
# =====================================================================
s9 = new_slide("08 · WHY IT WINS")
add_text(s9, Inches(0.6), Inches(1.55), Inches(12), Inches(1.4), "Not another job board.", size=42, bold=True, color=WHITE)

diffs = [
    ("Personalisation everywhere", "Match score isn’t hidden on one page — it drives Home, Dashboard, Skill Match, Detail and Applicant ranking. One user, one story."),
    ("Both sides, one product", "Students ship applications and pros ship shortlists in the same app. Every action on one side creates a signal on the other."),
    ("Judges can break it", "Change your skills → the app re-personalises live. Upload a real PDF → owners see it in the pipeline. Move a status → applicant is notified."),
    ("Production-lean", "10 seeded opportunities, deterministic match, real object storage, real email API, 18 backend tests, 0 console errors."),
]
top = Inches(2.75)
h = Inches(0.95)
for i, (title, body) in enumerate(diffs):
    y = top + i * (h + Inches(0.08))
    dot = s9.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), y + Inches(0.28), Inches(0.35), Inches(0.35))
    dot.fill.background(); dot.line.color.rgb = CYAN; dot.line.width = Pt(1.2); dot.shadow.inherit = False
    add_text(s9, Inches(0.6), y + Inches(0.26), Inches(0.35), Inches(0.35), str(i + 1), size=11, bold=True, color=CYAN, align=PP_ALIGN.CENTER)
    add_text(s9, Inches(1.15), y + Inches(0.15), Inches(4), Inches(0.5), title, size=16, bold=True, color=WHITE)
    add_text(s9, Inches(5.4), y + Inches(0.2), Inches(7.3), Inches(0.85), body, size=12.5, color=MUTED, line_spacing=1.5)


# =====================================================================
# 10. IMPACT & FUTURE
# =====================================================================
s10 = new_slide("09 · WHAT’S NEXT")
add_text(s10, Inches(0.6), Inches(1.55), Inches(12), Inches(1.4), "From hackathon to platform.", size=40, bold=True, color=WHITE)

roadmap = [
    ("NOW", CYAN, ["Two-sided demo ready", "Real object storage + email", "Deterministic AI match", "Full application pipeline"]),
    ("NEXT 4 WEEKS", VIOLET, ["Session-based auth + Google OAuth", "Verified sender domain for Resend", "Recruiter messaging inline", "Public shareable profile URL"]),
    ("NEXT 6 MONTHS", PINK, ["Interview scheduler", "Skill assessments (auto-verified)", "Team dashboards for colleges", "Freelance escrow (Stripe)"]),
]
top = Inches(3.15)
col_w = Inches(4.05)
gap = Inches(0.15)
for i, (title, color, items) in enumerate(roadmap):
    left = Inches(0.6) + i * (col_w + gap)
    add_rect(s10, left, top, col_w, Inches(3.6))
    acc = s10.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, col_w, Inches(0.06))
    acc.fill.solid(); acc.fill.fore_color.rgb = color; acc.line.fill.background(); acc.shadow.inherit = False
    add_eyebrow(s10, left + Inches(0.3), top + Inches(0.3), title, color=color, width=col_w)
    for j, item in enumerate(items):
        y = top + Inches(0.85) + j * Inches(0.55)
        chk = s10.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.3), y + Inches(0.05), Inches(0.18), Inches(0.18))
        chk.fill.solid(); chk.fill.fore_color.rgb = color; chk.line.fill.background(); chk.shadow.inherit = False
        add_text(s10, left + Inches(0.6), y, col_w - Inches(0.7), Inches(0.5), item, size=12.5, color=WHITE, line_spacing=1.4)


# =====================================================================
# 11. TEAM
# =====================================================================
s11 = new_slide("10 · THE TEAM")
add_text(s11, Inches(0.6), Inches(1.55), Inches(12), Inches(1.4), "Small team. Full stack. One sprint.", size=40, bold=True, color=WHITE)

team = [
    ("BUILDER 1", "Full-stack lead — React, FastAPI, Mongo, tests"),
    ("BUILDER 2", "AI + product logic — match engine, personalisation"),
    ("BUILDER 3", "Design & motion — dark-mode system, animations"),
    ("BUILDER 4", "Integrations — Resend email, object storage, deploys"),
]
top = Inches(3.15)
card_w = (SW - Inches(0.6) * 2 - Inches(0.15) * 3) / 4
for i, (role, blurb) in enumerate(team):
    left = Inches(0.6) + i * (card_w + Inches(0.15))
    add_rect(s11, left, top, card_w, Inches(3.4))
    # avatar
    av = s11.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.35), top + Inches(0.5), Inches(1.6), Inches(1.6))
    av.fill.solid(); av.fill.fore_color.rgb = [VIOLET, CYAN, PINK, AMBER][i]
    av.line.fill.background(); av.shadow.inherit = False
    add_text(s11, left + Inches(0.35), top + Inches(1.0), Inches(1.6), Inches(0.6), ["B1", "B2", "B3", "B4"][i], size=28, bold=True, color=BG, align=PP_ALIGN.CENTER)
    add_text(s11, left + Inches(0.3), top + Inches(2.35), card_w - Inches(0.6), Inches(0.4), role, size=11, bold=True, color=WHITE, font="Consolas", letter_spacing=200)
    add_text(s11, left + Inches(0.3), top + Inches(2.7), card_w - Inches(0.6), Inches(0.9), blurb, size=11, color=MUTED, line_spacing=1.5)

add_text(s11, Inches(0.6), Inches(6.75), Inches(12), Inches(0.4), "Replace with your teammates’ names, roles and photos before presenting.", size=11, color=MUTED, font="Consolas", letter_spacing=140)


# =====================================================================
# 12. THANK YOU / CTA
# =====================================================================
s12 = prs.slides.add_slide(BLANK)
add_bg(s12)
add_grid_bg(s12, size=Inches(0.85))
for cx, cy, rgb in [(Inches(9.5), Inches(1.8), VIOLET), (Inches(10.3), Inches(5.5), CYAN)]:
    blob = s12.shapes.add_shape(MSO_SHAPE.OVAL, cx, cy, Inches(3.8), Inches(3.8))
    blob.fill.solid(); blob.fill.fore_color.rgb = rgb
    blob.line.fill.background(); blob.fill.transparency = 0.85
    blob.shadow.inherit = False

add_brand_lockup(s12, Inches(0.6), Inches(0.45))
add_eyebrow(s12, Inches(0.6), Inches(2.4), "TRY IT · ASK US · JUDGE US", color=CYAN)

tb = s12.shapes.add_textbox(Inches(0.6), Inches(2.75), Inches(9.4), Inches(3.3))
tf = tb.text_frame
tf.word_wrap = True
p1 = tf.paragraphs[0]; p1.line_spacing = 1.02
r = p1.add_run(); r.text = "Bridge your skills."; r.font.name = "Calibri"; r.font.size = Pt(72); r.font.bold = True; r.font.color.rgb = WHITE
p2 = tf.add_paragraph(); p2.line_spacing = 1.02
r = p2.add_run(); r.text = "We built the bridge. "; r.font.name = "Calibri"; r.font.size = Pt(48); r.font.bold = True; r.font.color.rgb = WHITE
r = p2.add_run(); r.text = "Ready to walk it?"; r.font.name = "Calibri"; r.font.size = Pt(48); r.font.bold = True; r.font.color.rgb = VIOLET

add_text(s12, Inches(0.6), Inches(5.9), Inches(7.5), Inches(0.5), "Live demo · Student aarav@demo.com · Professional maya@demo.com", size=15, color=MUTED)

# CTA pill row
add_pill(s12, Inches(0.6), Inches(6.6), "TRY THE DEMO", CYAN, width=Inches(2.2))
add_pill(s12, Inches(3.0), Inches(6.6), "QUESTIONS · WELCOME", VIOLET, width=Inches(2.8))
add_pill(s12, Inches(6.0), Inches(6.6), "THANK YOU · JUDGES", PINK, width=Inches(2.4))

# QR card — top right
qr_box_left = Inches(10.15)
qr_box_top = Inches(2.55)
qr_box_w = Inches(2.7)
qr_box_h = Inches(3.9)
add_rect(s12, qr_box_left, qr_box_top, qr_box_w, qr_box_h, fill=RGBColor(0x11, 0x15, 0x22), line=CYAN)
add_eyebrow(s12, qr_box_left + Inches(0.25), qr_box_top + Inches(0.25), "SCAN · OPEN ON YOUR PHONE", color=CYAN, width=qr_box_w - Inches(0.5))
# White QR container (QR needs light background to scan well)
qr_pad_left = qr_box_left + Inches(0.35)
qr_pad_top = qr_box_top + Inches(0.7)
qr_pad_size = Inches(2.0)
white = s12.shapes.add_shape(MSO_SHAPE.RECTANGLE, qr_pad_left, qr_pad_top, qr_pad_size, qr_pad_size)
white.fill.solid(); white.fill.fore_color.rgb = RGBColor(0xF6, 0xF7, 0xFB); white.line.fill.background(); white.shadow.inherit = False
s12.shapes.add_picture(QR_PATH, qr_pad_left + Inches(0.1), qr_pad_top + Inches(0.1), width=qr_pad_size - Inches(0.2), height=qr_pad_size - Inches(0.2))
add_text(s12, qr_box_left + Inches(0.25), qr_pad_top + qr_pad_size + Inches(0.15), qr_box_w - Inches(0.5), Inches(0.5), "Try it while we talk.", size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s12, qr_box_left + Inches(0.25), qr_pad_top + qr_pad_size + Inches(0.5), qr_box_w - Inches(0.5), Inches(0.5), DEMO_URL.replace("https://", ""), size=9, color=CYAN, font="Consolas", align=PP_ALIGN.CENTER, letter_spacing=100)

slide_number(s12, 12, TOTAL)


out = "/app/SKILLBRIDGE_Hackathon_Deck.pptx"
prs.save(out)
print("Saved:", out)
